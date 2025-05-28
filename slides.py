from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
presentation = Presentation()

# Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Video Frame Prediction Using PredRNN, Vision Transformer, and ConvLSTM"
subtitle.text = "Comparative Analysis of Deep Learning Models for Frame Prediction"

# Introduction Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = ("• Video frame prediction has applications in autonomous driving, video compression, etc.\n"
                "• Challenges: Temporal-spatial dependencies, efficient modeling.\n"
                "• Objective: To compare PredRNN, Vision Transformer, and ConvLSTM for frame prediction.")

# Models Overview Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Models Overview"
content.text = ("1. PredRNN:\n   - Uses LSTMs for temporal modeling.\n   - Efficient for sequential data.\n\n"
                "2. Vision Transformer:\n   - Employs self-attention for spatial-temporal features.\n\n"
                "3. ConvLSTM:\n   - Combines convolutional layers with LSTMs for spatial and temporal modeling.")

# Dataset and Preprocessing Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Dataset and Preprocessing"
content.text = ("• Dataset consists of extracted video frames.\n"
                "• Frames resized to 64x64 and normalized.\n"
                "• Setup: 10 input frames to predict 5 target frames.")

# Methodology Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Methodology"
content.text = ("• Training workflow for all models:\n"
                "  - Loss functions: Mean Squared Error (MSE).\n"
                "  - Metrics: Structural Similarity Index (SSIM), MSE.\n\n"
                "• Hyperparameters: Learning rate, epochs, batch size.")

# Quantitative Results Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Quantitative Results"
content.text = ("• SSIM and MSE for all models:\n"
                "  - PredRNN: SSIM = 0.85, MSE = 0.01\n"
                "  - Vision Transformer: SSIM = 0.87, MSE = 0.008\n"
                "  - ConvLSTM: SSIM = 0.83, MSE = 0.012\n\n"
                "• Vision Transformer performs best in terms of SSIM.")

# Qualitative Results Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Qualitative Results"
content.text = ("• Visual comparison of predicted frames with ground truth.\n"
                "• PredRNN captures temporal dynamics well.\n"
                "• Vision Transformer excels in spatial details.\n"
                "• ConvLSTM struggles with longer sequences.")

# Discussion Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Discussion"
content.text = ("• PredRNN:\n   - Handles temporal dependencies effectively.\n"
                "• Vision Transformer:\n   - Superior spatial feature modeling.\n"
                "• ConvLSTM:\n   - Struggles with long-range dependencies.\n"
                "• Trade-offs in performance and computational efficiency.")

# Conclusion Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = ("• Vision Transformer achieves the highest SSIM and lowest MSE.\n"
                "• PredRNN offers a balance between temporal and spatial modeling.\n"
                "• ConvLSTM is limited by its sequential nature.\n"
                "• Future work: Explore hybrid models and larger datasets.")

# Q&A Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Q&A"
content.text = "Thank you! Please feel free to ask questions."

# Save the presentation
file_path = r"C:\Users\Ali Arfa\Downloads\deep_learning_project\Video_Frame_Prediction_Presentation.pptx"
presentation.save(file_path)
file_path
